from pptx import Presentation
import sys
import os

def find_first_pptx():
    for file in os.listdir('.'):
        if file.lower().endswith('.pptx'):
            return file
    return None

# Get .pptx file
if len(sys.argv) > 1:
    pptx_path = sys.argv[1]
else:
    pptx_path = find_first_pptx()
    if pptx_path is None:
        print("❌ No .pptx file found in current directory and no filename given.")
        sys.exit(1)

prs = Presentation(pptx_path)

print(f"==== SLIDE LAYOUTS IN '{pptx_path}' ====")
for idx, layout in enumerate(prs.slide_layouts):
    print(f"\nLayout {idx}: {layout.name}")
    print("  Placeholders:")
    for ph in layout.placeholders:
        pf = ph.placeholder_format
        print(f"    idx={pf.idx}, type={pf.type}, name='{ph.name}'")
        if hasattr(ph, "text"):
            print(f"      Text: '{ph.text.strip()}'")
    print("  Shapes:")
    for shape in layout.shapes:
        print(f"    name='{shape.name}', is_placeholder={shape.is_placeholder}, shape_type={shape.shape_type}")

print("\n==== SLIDES IN THIS FILE ====")
for idx, slide in enumerate(prs.slides):
    layout = slide.slide_layout
    print(f"\nSlide {idx + 1}: layout='{layout.name}'")
    print("  Shapes:")
    for shape in slide.shapes:
        print(f"    name='{shape.name}', is_placeholder={shape.is_placeholder}, shape_type={shape.shape_type}")
        if shape.is_placeholder:
            pf = shape.placeholder_format
            print(f"      Placeholder idx={pf.idx}, type={pf.type}")
        if hasattr(shape, "text"):
            txt = shape.text.strip()
            if txt:
                print(f"      Text: '{txt}'")
        if shape.has_table:
            print("      Contains table")
        if shape.has_chart:
            print("      Contains chart")
        if getattr(shape, "image", None) is not None:
            print("      Contains image")

print("\n==== END ====")
