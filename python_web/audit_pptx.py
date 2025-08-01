"""
audit_pptx.py - PowerPoint Layout and Slide Audit Utility

Purpose:
    - Audit the structure of a .pptx PowerPoint file.
    - Extracts and prints all slide layouts and their placeholders/shapes in Markdown table format.
    - Identifies each placeholder (idx, type, name, text) and all non-placeholder shapes (name, type, text, alt_text).
    - Helps technical users and developers understand the master layouts for slide automation or templating.

Usage:
    python audit_pptx.py [presentation.pptx]

Features:
    - Outputs Markdown tables for layouts and their elements.
    - Optionally saves the audit as 'filename_audit.md'.
    - Type numbers are taken directly from pptx; no mapping or naming applied.
    - Ensures alignment/padding in Markdown tables for human readability.

Editing/Customisation:
    - Output can be further processed for scripting slide generation, content mapping, or QA.

Limitations:
    - Designed for technical audit only; not for editing .pptx files.
    - Only works with files readable by python-pptx.

Author: [Your Name or Org]
"""

from pptx import Presentation
import sys
import os
import re

def find_first_pptx():
    for file in os.listdir('.'):
        if file.lower().endswith('.pptx'):
            return file
    return None

def pad_markdown_table(rows):
    if not rows:
        return []
    widths = [max(len(str(row[i])) for row in rows) for i in range(len(rows[0]))]
    return ["| " + " | ".join(str(cell).ljust(widths[i]) for i, cell in enumerate(row)) + " |" for row in rows]

def parse_block(name):
    """
    Parse the { - Key:Value ... } block in name.
    Returns a dict of keys/values.
    """
    block = {}
    if isinstance(name, str) and name.strip().startswith("{") and "-" in name:
        # Remove { and }, split lines, process "- key: value"
        lines = name.strip().strip("{}").splitlines()
        for line in lines:
            m = re.match(r"^\s*-\s*([\w\-]+)\s*:\s*(.*)", line)
            if m:
                key = m.group(1).strip()
                value = m.group(2).strip()
                block[key] = value
        return block if block else None
    return None

def collect_extra_cols(elements):
    # Gather all unique keys from block dicts
    cols = set()
    for elem in elements:
        if isinstance(elem, dict):
            cols.update(elem.keys())
    return sorted(cols)

def process_placeholders(placeholders_list):
    # Return (has_blocks, list of dicts, all block keys)
    has_block = False
    processed = []
    block_cols = set()
    for ph in placeholders_list:
        pf = ph.placeholder_format
        name = ph.name
        block = parse_block(name)
        if block:
            has_block = True
            # Still add idx/type for reference (optional: remove if not needed)
            block['idx'] = pf.idx
            block['type'] = pf.type
            block_cols.update(block.keys())
            processed.append(block)
        else:
            processed.append({
                "idx": pf.idx,
                "type": pf.type,
                "name": name,
                "text": ph.text.strip().replace('\n', ' ') if hasattr(ph, "text") and ph.text.strip() else ""
            })
    # Remove 'idx','type' from block_cols if only for reference
    block_cols = [col for col in block_cols if col not in ('idx', 'type')]
    return has_block, processed, block_cols

def process_shapes(shapes_list):
    has_block = False
    processed = []
    block_cols = set()
    for shape in shapes_list:
        if shape.is_placeholder:
            continue
        name = shape.name
        block = parse_block(name)
        if block:
            has_block = True
            block_cols.update(block.keys())
            processed.append(block)
        else:
            text = shape.text.strip().replace('\n', ' ') if hasattr(shape, "text") and shape.text.strip() else ""
            alt_text = shape.alternative_text.strip() if hasattr(shape, "alternative_text") else ""
            processed.append({
                "name": name,
                "shape_type": shape.shape_type,
                "text": text,
                "alt_text": alt_text
            })
    return has_block, processed, block_cols

if len(sys.argv) > 1:
    pptx_path = sys.argv[1]
else:
    pptx_path = find_first_pptx()
    if pptx_path is None:
        sys.exit(1)

prs = Presentation(pptx_path)
md_lines = []
md_lines.append(f"# Audit of '{pptx_path}'\n")
md_lines.append("## Slide Layouts\n\n")

for idx, layout in enumerate(prs.slide_layouts):
    md_lines.append(f"---\n\n### Layout {idx}: {layout.name}\n")
    # --- Placeholders
    phs = list(layout.placeholders)
    has_block, processed_phs, block_cols = process_placeholders(phs)
    if phs:
        if has_block:
            # Table header = all block keys found (sorted), idx/type always first if present
            headers = []
            if any('idx' in p for p in processed_phs): headers.append('idx')
            if any('type' in p for p in processed_phs): headers.append('type')
            headers += block_cols
            divider = ["---"] * len(headers)
            rows = []
            for ph in processed_phs:
                row = []
                for col in headers:
                    # error if not present at all
                    if col in ph:
                        row.append(ph[col])
                    else:
                        row.append("error")
                rows.append(row)
            padded = pad_markdown_table([headers, divider] + rows)
            md_lines.append("#### Placeholders\n")
            md_lines.extend(padded)
            md_lines.append("")
        else:
            header = ["idx", "type", "name", "text"]
            divider = ["---", "----", "----", "----"]
            rows = []
            for ph in processed_phs:
                row = [ph["idx"], ph["type"], ph["name"], ph["text"]]
                rows.append(row)
            padded = pad_markdown_table([header, divider] + rows)
            md_lines.append("#### Placeholders\n")
            md_lines.extend(padded)
            md_lines.append("")
    else:
        md_lines.append("#### Placeholders\n\n_None_\n")
    # --- Non-placeholder shapes
    sps = [shape for shape in layout.shapes if not shape.is_placeholder]
    has_block, processed_shapes, block_cols = process_shapes(sps)
    if sps:
        if has_block:
            headers = list(block_cols)
            divider = ["---"] * len(headers)
            rows = []
            for sp in processed_shapes:
                row = []
                for col in headers:
                    if col in sp:
                        row.append(sp[col])
                    else:
                        row.append("error")
                rows.append(row)
            padded = pad_markdown_table([headers, divider] + rows)
            md_lines.append("#### Non-placeholder shapes\n")
            md_lines.extend(padded)
            md_lines.append("")
        else:
            header = ["name", "shape_type", "text", "alt_text"]
            divider = ["----", "----------", "----", "--------"]
            rows = []
            for sp in processed_shapes:
                row = [sp["name"], sp["shape_type"], sp["text"], sp["alt_text"]]
                rows.append(row)
            padded = pad_markdown_table([header, divider] + rows)
            md_lines.append("#### Non-placeholder shapes\n")
            md_lines.extend(padded)
            md_lines.append("")
    else:
        md_lines.append("#### Non-placeholder shapes\n\n_None_\n")

md_lines.append("\n\n## Slides\n")

for idx, slide in enumerate(prs.slides):
    layout = slide.slide_layout
    md_lines.append(f"---\n\n### Slide {idx + 1}: layout='{layout.name}'\n")
    # --- Placeholders
    phs = [shape for shape in slide.shapes if shape.is_placeholder]
    has_block, processed_phs, block_cols = process_placeholders(phs)
    if phs:
        if has_block:
            headers = []
            if any('idx' in p for p in processed_phs): headers.append('idx')
            if any('type' in p for p in processed_phs): headers.append('type')
            headers += block_cols
            divider = ["---"] * len(headers)
            rows = []
            for ph in processed_phs:
                row = []
                for col in headers:
                    if col in ph:
                        row.append(ph[col])
                    else:
                        row.append("error")
                rows.append(row)
            padded = pad_markdown_table([headers, divider] + rows)
            md_lines.append("#### Placeholders\n")
            md_lines.extend(padded)
            md_lines.append("")
        else:
            header = ["idx", "type", "name", "text"]
            divider = ["---", "----", "----", "----"]
            rows = []
            for ph in processed_phs:
                row = [ph["idx"], ph["type"], ph["name"], ph["text"]]
                rows.append(row)
            padded = pad_markdown_table([header, divider] + rows)
            md_lines.append("#### Placeholders\n")
            md_lines.extend(padded)
            md_lines.append("")
    else:
        md_lines.append("#### Placeholders\n\n_None_\n")
    # --- Non-placeholder shapes
    sps = [shape for shape in slide.shapes if not shape.is_placeholder]
    has_block, processed_shapes, block_cols = process_shapes(sps)
    if sps:
        if has_block:
            headers = list(block_cols)
            divider = ["---"] * len(headers)
            rows = []
            for sp in processed_shapes:
                row = []
                for col in headers:
                    if col in sp:
                        row.append(sp[col])
                    else:
                        row.append("error")
                rows.append(row)
            padded = pad_markdown_table([headers, divider] + rows)
            md_lines.append("#### Non-placeholder shapes\n")
            md_lines.extend(padded)
            md_lines.append("")
        else:
            header = ["name", "shape_type", "text", "alt_text"]
            divider = ["----", "----------", "----", "--------"]
            rows = []
            for sp in processed_shapes:
                row = [sp["name"], sp["shape_type"], sp["text"], sp["alt_text"]]
                rows.append(row)
            padded = pad_markdown_table([header, divider] + rows)
            md_lines.append("#### Non-placeholder shapes\n")
            md_lines.extend(padded)
            md_lines.append("")
    else:
        md_lines.append("#### Non-placeholder shapes\n\n_None_\n")

md_output = "\n".join(md_lines)
print(md_output)

try:
    save = input("\nSave audit as markdown file (y/N)? ").strip().lower()
except Exception:
    save = 'n'
if save == "y":
    outname = os.path.splitext(os.path.basename(pptx_path))[0] + "_audit.md"
    with open(outname, "w") as f:
        f.write(md_output)
