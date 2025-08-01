import sys
import json
import os
from pptx import Presentation
import datetime

def find_first_file(ext):
    for file in os.listdir('.'):
        if file.lower().endswith(ext):
            return file
    return None

def get_timestamped_name(template):
    dt = datetime.datetime.now().strftime("%Y%m%d-%H%M%S")
    basename = os.path.splitext(template)[0]
    return f"{basename}-{dt}.pptx"

def main():
    # Allow optional arguments for .pptx and .json; else auto-find
    if len(sys.argv) == 3:
        pptx_file = sys.argv[1]
        input_file = sys.argv[2]
    else:
        pptx_file = find_first_file('.pptx')
        input_file = find_first_file('.json')
        if not pptx_file or not input_file:
            print("Usage: python3 pptx_gen.py <pptxfile> <input.json>")
            print("Or place a .pptx and .json file in the current directory.")
            sys.exit(1)

    prs = Presentation(pptx_file)
    with open(input_file, 'r') as f:
        slides = json.load(f)

    for slide_data in slides:
        layout_idx = slide_data.get("slidetype", 6)
        if layout_idx >= len(prs.slide_layouts):
            print(f"Warning: layout {layout_idx} not found. Skipping slide: {slide_data.get('title', '')}")
            continue
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        # Set title if present and placeholder available
        if "title" in slide_data and slide.placeholders:
            try:
                slide.placeholders[0].text = slide_data["title"]
            except Exception:
                pass
        # Set subtitle (placeholder type=4 for SUBTITLE)
        if "subtitle" in slide_data:
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == 4:
                    shape.text = slide_data["subtitle"]
        # Set bullets (placeholder idx=1 for Standard Text Slide)
        if "bullets" in slide_data and slide_data["bullets"] and len(slide.placeholders) > 1:
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            for idx, bullet in enumerate(slide_data["bullets"]):
                if idx == 0:
                    text_frame.text = bullet
                else:
                    p = text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 0

    out_file = get_timestamped_name(os.path.basename(pptx_file))
    prs.save(out_file)
    print(f"✅ Presentation saved as: {out_file}")

if __name__ == "__main__":
    main()
